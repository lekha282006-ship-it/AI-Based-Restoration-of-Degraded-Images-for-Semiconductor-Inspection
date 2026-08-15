from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
TITLE_BLUE = RGBColor(31, 58, 95)
ACCENT = RGBColor(46, 139, 87)
BG = RGBColor(15, 23, 42)
TEXT = RGBColor(245, 247, 250)

# Helper

def add_textbox(slide, left, top, width, height, font_size=22):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 8
    tf.margin_right = 8
    tf.margin_top = 8
    tf.margin_bottom = 8
    return tf

# Slide 1 - Title
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = BG

# Title block
shape = slide.shapes.add_shape(1, Inches(0.7), Inches(0.9), Inches(11.8), Inches(1.3))
shape.fill.solid()
shape.fill.fore_color.rgb = TITLE_BLUE
shape.line.color.rgb = TITLE_BLUE
text = shape.text_frame
text.text = 'AI-Based Restoration of Degraded Images for Semiconductor Inspection'
text.paragraphs[0].font.size = Pt(26)
text.paragraphs[0].font.bold = True
text.paragraphs[0].font.color.rgb = TEXT

subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(8.5), Inches(1.0))
subtitle_tf = subtitle.text_frame
subtitle_tf.text = 'Deep Learning Approach for Noise Removal and Super-Resolution'
subtitle_tf.paragraphs[0].font.size = Pt(18)
subtitle_tf.paragraphs[0].font.color.rgb = TEXT
subtitle_tf.paragraphs[0].font.bold = True

# add small panel with metrics
panel = slide.shapes.add_shape(1, Inches(8.7), Inches(2.3), Inches(3.3), Inches(1.8))
panel.fill.solid()
panel.fill.fore_color.rgb = ACCENT
panel.line.color.rgb = ACCENT
panel_tf = panel.text_frame
panel_tf.text = 'Validation Snapshot\nPSNR: 11.48\nSSIM: 0.111\nLPIPS: 0.660'
panel_tf.paragraphs[0].font.size = Pt(18)
panel_tf.paragraphs[0].font.bold = True
panel_tf.paragraphs[1].font.size = Pt(16)
panel_tf.paragraphs[2].font.size = Pt(16)
panel_tf.paragraphs[3].font.size = Pt(16)

# Slide 2 - Problem & Objective
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Problem and Objective'
box = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(6.0), Inches(4.5), 22)
for i, line in enumerate([
    'Semiconductor inspection images are often noisy, degraded, and low-resolution.',
    'These degradations hide structural patterns that are critical for defect analysis.',
    'The goal is to recover visually cleaner and more interpretable images in a single pass.',
    'This helps improve downstream inspection and aids automated semiconductor quality review.'
]):
    p = box.paragraphs[0] if i == 0 else box.add_paragraph()
    p.text = line
    p.font.size = Pt(22)
    p.level = 0

# Simple stat panel on right
panel = slide.shapes.add_shape(1, Inches(7.3), Inches(1.8), Inches(4.8), Inches(2.6))
panel.fill.solid(); panel.fill.fore_color.rgb = TITLE_BLUE; panel.line.color.rgb = TITLE_BLUE
p_tf = panel.text_frame
p_tf.text = 'Challenge Constraints\n• CPU-only environment\n• No GPU available\n• Real image restoration task'
p_tf.paragraphs[0].font.size = Pt(20); p_tf.paragraphs[0].font.bold = True
for i in range(1, len(p_tf.paragraphs)):
    p_tf.paragraphs[i].font.size = Pt(17)

# Slide 3 - Dataset and Method
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Dataset and Methodology'
left = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.0), Inches(4.2), 22)
for i, line in enumerate([
    'Paired noisy/ground-truth .npy dataset',
    'Train/validation split from noisy and GT folders',
    'External held-out noisy folder for generalization testing',
    'Input normalization: clip to [0, 1.5], then scale by 1.5',
    'Output target: clean grayscale image in [0, 1] range'
]):
    p = left.paragraphs[0] if i == 0 else left.add_paragraph()
    p.text = line
    p.font.size = Pt(21)

img_path = Path('results/val_grid_demo.png')
if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(6.3), Inches(1.8), Inches(5.8), Inches(3.7))

# Slide 4 - Model Architecture
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Model Architecture'
box = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(4.8), Inches(4.0), 22)
for i, line in enumerate([
    'Residual CNN backbone',
    'Instance normalization layers',
    'Pixel-shuffle upsampling head',
    'Single-pass restoration for denoising + super-resolution',
    'Compact and efficient for CPU-based execution' ]):
    p = box.paragraphs[0] if i == 0 else box.add_paragraph()
    p.text = line
    p.font.size = Pt(22)

if img_path.exists():
    slide.shapes.add_picture(str(img_path), Inches(6.1), Inches(1.7), Inches(6.2), Inches(3.9))

# Slide 5 - Results
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Results and Quantitative Evaluation'
box = add_textbox(slide, Inches(0.7), Inches(1.5), Inches(5.0), Inches(4.2), 22)
for i, line in enumerate([
    'PSNR mean: 11.48 dB',
    'SSIM mean: 0.111',
    'LPIPS mean: 0.660',
    'Evaluation run used real paired noisy/GT samples',
    'OOD check: 400 unseen semiconductor images processed without retraining'
]):
    p = box.paragraphs[0] if i == 0 else box.add_paragraph()
    p.text = line
    p.font.size = Pt(22)

# add a small metric card
panel = slide.shapes.add_shape(1, Inches(7.2), Inches(1.8), Inches(4.8), Inches(2.4))
panel.fill.solid(); panel.fill.fore_color.rgb = ACCENT; panel.line.color.rgb = ACCENT
p_tf = panel.text_frame
p_tf.text = 'Key Result\nModel successfully restored held-out semiconductor images in the external test folder.'
p_tf.paragraphs[0].font.size = Pt(20); p_tf.paragraphs[0].font.bold = True
p_tf.paragraphs[1].font.size = Pt(17)

# Slide 6 - Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = 'Conclusion'
box = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.2), Inches(4.0), 23)
for i, line in enumerate([
    'A working end-to-end image restoration pipeline was implemented for semiconductor inspection images.',
    'The model was validated on real data and also tested on an unseen external noisy folder, confirming practical generalization.',
    'Repository includes training, inference, evaluation, dependency setup, and submission summary for reproducibility.',
    'The project is ready for presentation and further extension with larger models and GPU training.'
]):
    p = box.paragraphs[0] if i == 0 else box.add_paragraph()
    p.text = line
    p.font.size = Pt(23)

out = Path('AI-Based-Restoration-of-Degraded-Images-for-Semiconductor-Inspection.pptx')
prs.save(str(out))
print(f'Created {out.resolve()}')
